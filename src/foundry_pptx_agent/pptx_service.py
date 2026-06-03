from __future__ import annotations

import json
import re
import time
import zipfile
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER
from pptx.util import Inches, Pt

from .schemas import CreateDeckRequest, CreateDeckResponse, DeckPlan, SlidePlan, as_posix_str
from .settings import OUTPUT_DIR, artifact_url, template_contract_path, template_path

PLACEHOLDER_PATTERNS = [
    re.compile(pattern, re.IGNORECASE)
    for pattern in [
        r"\{\{[^}]+\}\}",
        r"\bxxxx+\b",
        r"\blorem\b",
        r"\bipsum\b",
        r"click to add",
    ]
]


def list_templates() -> list[dict[str, str]]:
    templates = []
    for path in sorted(template_path.parent.glob("*.pptx")):
        templates.append({"template_id": path.stem, "path": as_posix_str(path)})
    return templates


def analyze_template(template_id: str) -> dict[str, Any]:
    path = template_path(template_id)
    prs = Presentation(path)
    layouts = []
    for index, layout in enumerate(prs.slide_layouts):
        placeholders = []
        for shape in layout.placeholders:
            placeholders.append(
                {
                    "idx": shape.placeholder_format.idx,
                    "name": shape.name,
                    "type": str(shape.placeholder_format.type),
                    "left": int(shape.left),
                    "top": int(shape.top),
                    "width": int(shape.width),
                    "height": int(shape.height),
                }
            )
        layouts.append(
            {
                "index": index,
                "name": layout.name,
                "placeholder_count": len(placeholders),
                "placeholders": placeholders,
            }
        )
    return {
        "template_id": template_id,
        "path": as_posix_str(path),
        "slide_width": int(prs.slide_width),
        "slide_height": int(prs.slide_height),
        "layout_count": len(layouts),
        "layouts": layouts,
        "contract": _load_contract(template_id),
    }


def create_deck(request: CreateDeckRequest) -> CreateDeckResponse:
    template_id = request.template_id or request.deck_plan.deck.template_id
    template = template_path(template_id)
    prs = Presentation(template)
    _remove_existing_slides(prs)
    contract = _load_contract(template_id)
    brand = _brand(contract)

    slide_mapping: list[dict[str, Any]] = []
    for slide_plan in request.deck_plan.slides:
        layout_index, reason = _resolve_layout(prs, slide_plan, contract)
        slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
        if contract.get("render_mode") == "master_layout":
            _apply_master_layout_content(slide, slide_plan, contract)
        else:
            _apply_slide_design(prs, slide, slide_plan, brand)
        slide_mapping.append(
            {
                "slide_number": slide_plan.slide_number,
                "intent": slide_plan.intent,
                "title": slide_plan.title,
                "layout_index": layout_index,
                "layout_name": prs.slide_layouts[layout_index].name,
                "reason": reason,
            }
        )

    output_name = request.output_name or _slug(request.deck_plan.deck.title)
    if not output_name.endswith(".pptx"):
        output_name = f"{output_name}.pptx"
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    output_path = (OUTPUT_DIR / output_name).resolve()
    prs.save(output_path)
    validation = validate_deck(output_path, template_id)
    return CreateDeckResponse(
        artifact_path=as_posix_str(output_path),
        artifact_url=artifact_url(output_path.name),
        template_id=template_id,
        slide_count=len(request.deck_plan.slides),
        validation=validation,
        slide_mapping=slide_mapping,
    )


def validate_deck(deck_path: str | Path, template_id: str = "sample-board-template") -> dict[str, Any]:
    path = Path(deck_path).resolve()
    errors: list[str] = []
    warnings: list[str] = []
    placeholder_hits: list[dict[str, str]] = []

    if path.suffix.lower() == ".pptm":
        errors.append("Macro-enabled .pptm files are not supported by this prototype.")
    if not path.exists():
        errors.append(f"Deck not found: {path}")
        return {"ok": False, "errors": errors, "warnings": warnings}
    if not zipfile.is_zipfile(path):
        errors.append("Deck is not a valid zip-based PPTX package.")
        return {"ok": False, "errors": errors, "warnings": warnings}

    try:
        prs = Presentation(path)
    except Exception as exc:
        errors.append(f"PowerPoint package could not be opened: {exc}")
        return {"ok": False, "errors": errors, "warnings": warnings}

    contract = _load_contract(template_id)
    limits = contract.get("limits", {})
    for index, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = "\n".join(paragraph.text for paragraph in shape.text_frame.paragraphs).strip()
            if not text:
                continue
            for pattern in PLACEHOLDER_PATTERNS:
                if pattern.search(text):
                    placeholder_hits.append({"slide": str(index), "text": text[:120]})
            title_limit = int(limits.get("title_chars", 72))
            body_limit = int(limits.get("message_chars", 220))
            if len(text) > body_limit and not _looks_like_footer(text):
                warnings.append(f"Slide {index} has a long text block ({len(text)} chars).")
            if len(text) > title_limit and shape == _title_shape(slide):
                warnings.append(f"Slide {index} title may be too long ({len(text)} chars).")

    if placeholder_hits:
        errors.append("Placeholder or sample text remains in the generated deck.")

    return {
        "ok": not errors,
        "errors": errors,
        "warnings": warnings,
        "placeholder_hits": placeholder_hits,
        "slide_count": len(prs.slides),
        "validated_at": int(time.time()),
    }


def _load_contract(template_id: str) -> dict[str, Any]:
    path = template_contract_path(template_id)
    if not path.exists():
        return {}
    with path.open("r", encoding="utf-8") as handle:
        return json.load(handle)


def _brand(contract: dict[str, Any]) -> dict[str, str]:
    return {
        "primary": contract.get("brand", {}).get("primary", "1E2761"),
        "secondary": contract.get("brand", {}).get("secondary", "CADCFC"),
        "accent": contract.get("brand", {}).get("accent", "00A896"),
        "background": contract.get("brand", {}).get("background", "FFFFFF"),
        "text": contract.get("brand", {}).get("text", "1F2937"),
        "muted": contract.get("brand", {}).get("muted", "64748B"),
        "font_face": contract.get("brand", {}).get("font_face", "Aptos"),
        "footer": contract.get("required_footer", "Generated by Foundry PPTX Agent"),
    }


def _remove_existing_slides(prs: Presentation) -> None:
    slide_id_list = prs.slides._sldIdLst  # python-pptx has no public delete API.
    for slide_id in list(slide_id_list):
        rel_id = slide_id.rId
        prs.part.drop_rel(rel_id)
        slide_id_list.remove(slide_id)


def _resolve_layout(prs: Presentation, slide_plan: SlidePlan, contract: dict[str, Any] | None = None) -> tuple[int, str]:
    layout_map = (contract or {}).get("layout_map", {})
    configured = layout_map.get(slide_plan.intent) or layout_map.get(slide_plan.layout_hint)
    if configured is not None:
        if isinstance(configured, int) and 0 <= configured < len(prs.slide_layouts):
            return configured, "matched contract layout_map index"
        if isinstance(configured, str):
            for index, layout in enumerate(prs.slide_layouts):
                if configured.lower() == layout.name.lower():
                    return index, "matched contract layout_map name"

    names = [layout.name.lower() for layout in prs.slide_layouts]
    hint = slide_plan.layout_hint.lower().strip()
    intent = slide_plan.intent.lower().strip()
    for index, name in enumerate(names):
        if hint and hint in name:
            return index, "matched layout_hint to layout name"
    for token in _intent_tokens(intent):
        for index, name in enumerate(names):
            if token in name:
                return index, f"matched intent token '{token}' to layout name"
    for index, layout in enumerate(prs.slide_layouts):
        placeholder_types = {shape.placeholder_format.type for shape in layout.placeholders}
        if PP_PLACEHOLDER.TITLE in placeholder_types and PP_PLACEHOLDER.BODY in placeholder_types:
            return index, "matched title/body placeholder signature"
    return 0, "fell back to first layout"


def _intent_tokens(intent: str) -> list[str]:
    mapping = {
        "title": ["title"],
        "conclusion": ["title", "section"],
        "executive_summary": ["content", "summary"],
        "comparison": ["comparison", "two", "content"],
        "timeline": ["content", "section"],
        "risks": ["content"],
    }
    return mapping.get(intent, ["content", "title"])


def _apply_master_layout_content(slide: Any, plan: SlidePlan, contract: dict[str, Any]) -> None:
    text_placeholders = _text_placeholders(slide)
    if not text_placeholders:
        return

    title_shape = _title_shape(slide)
    if title_shape is not None:
        _write_text_frame(title_shape.text_frame, [plan.title])

    body_shapes = [shape for shape in text_placeholders if shape != title_shape]
    if not body_shapes:
        return

    lines = _content_lines(plan)
    if len(body_shapes) == 1:
        _write_text_frame(body_shapes[0].text_frame, lines)
        return

    blocks = plan.content_blocks
    columns = _placeholder_columns(body_shapes)
    if blocks and len(columns) >= len(blocks) and all(len(column) >= 2 for column in columns[: len(blocks)]):
        used_shapes: list[Any] = []
        for column, block in zip(columns, blocks, strict=False):
            _write_text_frame(column[0].text_frame, [block.label])
            _write_text_frame(column[1].text_frame, [block.text])
            used_shapes.extend(column[:2])
        for shape in body_shapes:
            if shape not in used_shapes:
                shape.text_frame.clear()
        return

    if blocks and len(body_shapes) >= len(blocks) * 2:
        paired = _top_then_body_pairs(body_shapes)
        for (label_shape, body_shape), block in zip(paired, blocks, strict=False):
            _write_text_frame(label_shape.text_frame, [block.label])
            _write_text_frame(body_shape.text_frame, [block.text])
        for shape in body_shapes[len(paired) * 2 :]:
            shape.text_frame.clear()
        return

    values = _distributed_values(plan, len(body_shapes))
    for shape, value in zip(body_shapes, values, strict=False):
        _write_text_frame(shape.text_frame, value)
    for shape in body_shapes[len(values) :]:
        shape.text_frame.clear()


def _text_placeholders(slide: Any) -> list[Any]:
    placeholders: list[Any] = []
    for shape in slide.placeholders:
        if not getattr(shape, "has_text_frame", False):
            continue
        try:
            placeholder_type = shape.placeholder_format.type
        except ValueError:
            continue
        if placeholder_type in {
            PP_PLACEHOLDER.TITLE,
            PP_PLACEHOLDER.CENTER_TITLE,
            PP_PLACEHOLDER.BODY,
            PP_PLACEHOLDER.OBJECT,
            PP_PLACEHOLDER.SUBTITLE,
        }:
            placeholders.append(shape)
    return placeholders


def _placeholder_idx(shape: Any) -> int:
    try:
        return int(shape.placeholder_format.idx)
    except ValueError:
        return 999


def _placeholder_columns(shapes: list[Any]) -> list[list[Any]]:
    columns: list[list[Any]] = []
    tolerance = 91440  # 0.1 inch in EMUs.
    for shape in sorted(shapes, key=lambda item: (int(item.left), int(item.top))):
        for column in columns:
            if abs(int(column[0].left) - int(shape.left)) <= tolerance:
                column.append(shape)
                break
        else:
            columns.append([shape])
    for column in columns:
        column.sort(key=lambda item: int(item.top))
    return columns


def _top_then_body_pairs(shapes: list[Any]) -> list[tuple[Any, Any]]:
    ordered = sorted(shapes, key=lambda item: (int(item.top), int(item.left)))
    midpoint = len(ordered) // 2
    labels = sorted(ordered[:midpoint], key=lambda item: int(item.left))
    bodies = sorted(ordered[midpoint:], key=lambda item: int(item.left))
    return list(zip(labels, bodies, strict=False))


def _content_lines(plan: SlidePlan) -> list[str]:
    lines = [plan.message] if plan.message else []
    for block in plan.content_blocks:
        if block.label:
            lines.append(block.label)
        lines.append(block.text)
    return [line for line in lines if line]


def _distributed_values(plan: SlidePlan, count: int) -> list[list[str]]:
    if not plan.content_blocks:
        return [[plan.message]] if plan.message else []
    values: list[list[str]] = []
    for block in plan.content_blocks[:count]:
        values.append([line for line in [block.label, block.text] if line])
    return values


def _write_text_frame(text_frame: Any, lines: list[str]) -> None:
    text_frame.clear()
    text_frame.word_wrap = True
    if not lines:
        return
    first = text_frame.paragraphs[0]
    first.text = lines[0]
    for line in lines[1:]:
        paragraph = text_frame.add_paragraph()
        paragraph.text = line


def _apply_slide_design(prs: Presentation, slide: Any, plan: SlidePlan, brand: dict[str, str]) -> None:
    if plan.intent in {"title", "conclusion"}:
        _render_title_slide(prs, slide, plan, brand)
    elif plan.intent == "comparison":
        _render_comparison_slide(prs, slide, plan, brand)
    elif plan.intent == "timeline":
        _render_timeline_slide(prs, slide, plan, brand)
    else:
        _render_cards_slide(prs, slide, plan, brand)
    _add_footer(prs, slide, brand)


def _render_title_slide(prs: Presentation, slide: Any, plan: SlidePlan, brand: dict[str, str]) -> None:
    _clear_placeholders(slide)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(brand["primary"])
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(4.6), prs.slide_width, Inches(1.05))
    accent.fill.solid()
    accent.fill.fore_color.rgb = _rgb(brand["accent"])
    accent.line.color.rgb = _rgb(brand["accent"])
    title_box = slide.shapes.add_textbox(Inches(0.65), Inches(1.1), Inches(8.7), Inches(1.2))
    _set_text(title_box, plan.title, 40, "FFFFFF", brand, bold=True)
    msg_box = slide.shapes.add_textbox(Inches(0.68), Inches(2.45), Inches(7.8), Inches(1.1))
    _set_text(msg_box, plan.message, 18, "CADCFC", brand)
    if plan.content_blocks:
        callout = plan.content_blocks[0]
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(3.6), Inches(7.8), Inches(0.7))
        box.fill.solid()
        box.fill.fore_color.rgb = _rgb("FFFFFF")
        box.line.color.rgb = _rgb("FFFFFF")
        _set_text(box, f"{callout.label}: {callout.text}", 14, brand["primary"], brand, bold=True)


def _render_cards_slide(prs: Presentation, slide: Any, plan: SlidePlan, brand: dict[str, str]) -> None:
    _clear_placeholders(slide)
    _add_header(prs, slide, plan, brand)
    message = slide.shapes.add_textbox(Inches(0.65), Inches(1.15), Inches(8.7), Inches(0.7))
    _set_text(message, plan.message, 15, brand["muted"], brand)
    blocks = plan.content_blocks or []
    count = max(1, min(len(blocks), 4))
    card_w = 8.8 / count
    for index, block in enumerate(blocks[:4]):
        x = 0.6 + index * card_w
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.15), Inches(card_w - 0.18), Inches(2.35))
        card.fill.solid()
        card.fill.fore_color.rgb = _rgb("F8FAFC")
        card.line.color.rgb = _rgb(brand["secondary"])
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(2.15), Inches(0.08), Inches(2.35))
        accent.fill.solid()
        accent.fill.fore_color.rgb = _rgb(brand["accent"])
        accent.line.color.rgb = _rgb(brand["accent"])
        label = slide.shapes.add_textbox(Inches(x + 0.22), Inches(2.38), Inches(card_w - 0.55), Inches(0.35))
        _set_text(label, block.label, 15, brand["primary"], brand, bold=True)
        body = slide.shapes.add_textbox(Inches(x + 0.22), Inches(2.9), Inches(card_w - 0.55), Inches(1.25))
        _set_text(body, block.text, 12, brand["text"], brand)


def _render_comparison_slide(prs: Presentation, slide: Any, plan: SlidePlan, brand: dict[str, str]) -> None:
    _clear_placeholders(slide)
    _add_header(prs, slide, plan, brand)
    blocks = plan.content_blocks[:2]
    for index, block in enumerate(blocks):
        x = 0.75 + index * 4.55
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.55), Inches(3.95), Inches(3.2))
        card.fill.solid()
        card.fill.fore_color.rgb = _rgb("FFFFFF" if index else "F8FAFC")
        card.line.color.rgb = _rgb(brand["accent"] if index else brand["secondary"])
        label = slide.shapes.add_textbox(Inches(x + 0.35), Inches(1.9), Inches(3.25), Inches(0.5))
        _set_text(label, block.label, 18, brand["primary"], brand, bold=True)
        body = slide.shapes.add_textbox(Inches(x + 0.35), Inches(2.65), Inches(3.25), Inches(1.45))
        _set_text(body, block.text, 14, brand["text"], brand)
    msg = slide.shapes.add_textbox(Inches(0.85), Inches(4.95), Inches(8.3), Inches(0.35))
    _set_text(msg, plan.message, 12, brand["muted"], brand)


def _render_timeline_slide(prs: Presentation, slide: Any, plan: SlidePlan, brand: dict[str, str]) -> None:
    _clear_placeholders(slide)
    _add_header(prs, slide, plan, brand)
    blocks = plan.content_blocks[:4]
    y = 2.65
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.1), Inches(y + 0.18), Inches(7.8), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = _rgb(brand["secondary"])
    line.line.color.rgb = _rgb(brand["secondary"])
    for index, block in enumerate(blocks):
        x = 0.85 + index * 2.25
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.5), Inches(0.5))
        dot.fill.solid()
        dot.fill.fore_color.rgb = _rgb(brand["accent"])
        dot.line.color.rgb = _rgb(brand["accent"])
        number = slide.shapes.add_textbox(Inches(x), Inches(y + 0.09), Inches(0.5), Inches(0.2))
        _set_text(number, str(index + 1), 11, "FFFFFF", brand, bold=True, align="center")
        label = slide.shapes.add_textbox(Inches(x - 0.25), Inches(y + 0.7), Inches(1.8), Inches(0.35))
        _set_text(label, block.label, 14, brand["primary"], brand, bold=True, align="center")
        body = slide.shapes.add_textbox(Inches(x - 0.45), Inches(y + 1.15), Inches(2.0), Inches(0.9))
        _set_text(body, block.text, 10.5, brand["text"], brand, align="center")
    msg = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(0.55))
    _set_text(msg, plan.message, 14, brand["muted"], brand)


def _add_header(prs: Presentation, slide: Any, plan: SlidePlan, brand: dict[str, str]) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(brand["primary"])
    bar.line.color.rgb = _rgb(brand["primary"])
    title = slide.shapes.add_textbox(Inches(0.6), Inches(0.45), Inches(8.8), Inches(0.55))
    _set_text(title, plan.title, 25, brand["primary"], brand, bold=True)


def _add_footer(prs: Presentation, slide: Any, brand: dict[str, str]) -> None:
    footer = slide.shapes.add_textbox(Inches(0.6), Inches(5.25), Inches(8.8), Inches(0.2))
    _set_text(footer, brand["footer"], 8, brand["muted"], brand)


def _clear_placeholders(slide: Any) -> None:
    for shape in list(slide.shapes):
        if getattr(shape, "is_placeholder", False) and getattr(shape, "has_text_frame", False):
            shape.text_frame.clear()


def _title_shape(slide: Any) -> Any | None:
    for shape in slide.shapes:
        if getattr(shape, "is_placeholder", False):
            if shape.placeholder_format.type in {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}:
                return shape
    return None


def _looks_like_footer(text: str) -> bool:
    return "generated by foundry pptx agent" in text.lower()


def _set_text(shape: Any, text: str, size: float, color: str, brand: dict[str, str], bold: bool = False, align: str | None = None) -> None:
    shape.text_frame.clear()
    shape.text_frame.word_wrap = True
    paragraph = shape.text_frame.paragraphs[0]
    paragraph.text = text
    if align == "center":
        paragraph.alignment = 2
    run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
    run.font.name = brand["font_face"]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)


def _rgb(hex_color: str) -> RGBColor:
    clean = hex_color.strip().lstrip("#")
    return RGBColor(int(clean[0:2], 16), int(clean[2:4], 16), int(clean[4:6], 16))


def _slug(value: str) -> str:
    slug = re.sub(r"[^a-zA-Z0-9]+", "-", value.lower()).strip("-")
    return slug or "generated-deck"
