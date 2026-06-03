from __future__ import annotations

import json
import re
import shutil
from pathlib import Path
from typing import Any

from pptx import Presentation
from .settings import TEMPLATES_DIR, template_contract_path


DEFAULT_INTENTS = [
    "title",
    "executive_summary",
    "comparison",
    "timeline",
    "risks",
    "conclusion",
]


def onboard_template(source_path: str | Path, template_id: str | None = None, overwrite: bool = False) -> dict[str, Any]:
    source = Path(source_path).expanduser().resolve()
    if not source.exists() or source.suffix.lower() != ".pptx":
        raise ValueError(f"Expected a .pptx file, got: {source}")

    resolved_template_id = _safe_template_id(template_id or source.stem)
    destination = template_path_for_write(resolved_template_id)
    contract_destination = template_contract_path(resolved_template_id)
    if not overwrite and (destination.exists() or contract_destination.exists()):
        raise FileExistsError(
            f"Template '{resolved_template_id}' already exists. Use overwrite=True to replace it."
        )

    TEMPLATES_DIR.mkdir(parents=True, exist_ok=True)
    shutil.copy2(source, destination)

    contract = build_contract(destination, resolved_template_id)
    with contract_destination.open("w", encoding="utf-8") as handle:
        json.dump(contract, handle, indent=2)
        handle.write("\n")

    return {
        "template_id": resolved_template_id,
        "template_path": str(destination),
        "contract_path": str(contract_destination),
        "contract": contract,
    }


def build_contract(template_file: str | Path, template_id: str) -> dict[str, Any]:
    prs = Presentation(template_file)
    layout_catalog = [_layout_entry(index, layout) for index, layout in enumerate(prs.slide_layouts)]
    layout_map = _infer_layout_map(layout_catalog)
    return {
        "template_id": template_id,
        "name": prs.core_properties.title or template_id,
        "version": "local",
        "render_mode": "master_layout",
        "layout_map": layout_map,
        "layout_catalog": layout_catalog,
        "limits": {
            "title_chars": 90,
            "message_chars": 260,
            "block_chars": 180,
        },
    }


def template_path_for_write(template_id: str) -> Path:
    safe_id = _safe_template_id(template_id)
    path = (TEMPLATES_DIR / f"{safe_id}.pptx").resolve()
    if not str(path).startswith(str(TEMPLATES_DIR)):
        raise ValueError("template_id resolved outside the templates directory")
    return path


def _layout_entry(index: int, layout: Any) -> dict[str, Any]:
    placeholders = []
    for shape in layout.placeholders:
        placeholders.append(
            {
                "idx": shape.placeholder_format.idx,
                "name": shape.name,
                "type": str(shape.placeholder_format.type),
                "left": int(shape.left),
                "top": int(shape.top),
                "width": int(shape.width),
                "height": int(shape.height),
            }
        )
    return {
        "index": index,
        "name": layout.name,
        "placeholder_count": len(placeholders),
        "text_placeholder_count": sum(1 for item in placeholders if "TITLE" in item["type"] or "BODY" in item["type"] or "OBJECT" in item["type"]),
        "picture_placeholder_count": sum(1 for item in placeholders if "PICTURE" in item["type"]),
        "placeholders": placeholders,
    }


def _infer_layout_map(layout_catalog: list[dict[str, Any]]) -> dict[str, int]:
    return {
        "title": _pick_layout(layout_catalog, ["title slide", "title"], min_text=2),
        "executive_summary": _pick_layout(layout_catalog, ["three column", "four column", "3_4-column", "title and content", "content"], min_text=3),
        "comparison": _pick_layout(layout_catalog, ["two column", "side by side", "comparison"], min_text=3),
        "timeline": _pick_layout(layout_catalog, ["four column", "five column", "3_4-column", "timeline", "roadmap"], min_text=5),
        "risks": _pick_layout(layout_catalog, ["three column", "four column", "3_4-column", "title and content", "content"], min_text=3),
        "conclusion": _pick_layout(layout_catalog, ["title slide 2", "closing", "section title", "title"], min_text=1),
    }


def _pick_layout(layout_catalog: list[dict[str, Any]], preferred_terms: list[str], min_text: int) -> int:
    scored: list[tuple[int, int, dict[str, Any]]] = []
    for entry in layout_catalog:
        name = entry["name"].lower()
        score = 0
        for offset, term in enumerate(preferred_terms):
            if term in name:
                score += (len(preferred_terms) - offset) * 100
        text_count = int(entry["text_placeholder_count"])
        if text_count >= min_text:
            score += 20
        score -= abs(text_count - min_text)
        if int(entry["picture_placeholder_count"]) > 0:
            score -= 5
        scored.append((score, text_count, entry))
    scored.sort(key=lambda item: (item[0], item[1]), reverse=True)
    return int(scored[0][2]["index"]) if scored else 0


def _safe_template_id(value: str) -> str:
    normalized = re.sub(r"[^a-zA-Z0-9_-]+", "-", value.strip().lower()).strip("-")
    if not normalized:
        raise ValueError("template_id cannot be empty")
    return normalized[:80]
