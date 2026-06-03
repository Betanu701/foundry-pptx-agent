from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "templates" / "sample-board-template.pptx"


def main() -> None:
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.core_properties.title = "Sample Board Template"
    prs.core_properties.subject = "Template fixture for Foundry PPTX Agent"
    prs.core_properties.author = "Foundry PPTX Agent"
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(0x1E, 0x27, 0x61)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(4.6), Inches(10), Inches(1.05))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(0x00, 0xA8, 0x96)
    accent.line.color.rgb = RGBColor(0x00, 0xA8, 0x96)
    title = slide.shapes.title
    title.text = "Board Template"
    for paragraph in title.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = "Aptos"
            run.font.size = Pt(40)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    subtitle = slide.placeholders[1]
    subtitle.text = "Customer-ready demo template"
    for paragraph in subtitle.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = "Aptos"
            run.font.size = Pt(18)
            run.font.color.rgb = RGBColor(0xCA, 0xDC, 0xFC)
    prs.save(OUTPUT)


if __name__ == "__main__":
    main()

